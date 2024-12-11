import os
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor

class PowerPointEditor:
    """
    A class to encapsulate PowerPoint text modification operations.
    """

    def __init__(self, input_file: str):
        self.input_file = input_file
        self.presentation = Presentation(input_file)

    def convert_text_to_uppercase(self):
        """
        Converts all text in the presentation to uppercase.
        """
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:  # Check if the shape contains text
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.upper()

    def save(self, output_file: str):
        """
        Saves the modified presentation to the output file.
        """
        self.presentation.save(output_file)
        print(f"Modified presentation saved to: {output_file}")


class PowerPointBatchProcessor:
    """
    A class to handle batch processing of PowerPoint files.
    """

    def __init__(self, input_folder: str, output_folder: str):
        self.input_folder = input_folder
        self.output_folder = output_folder
        os.makedirs(output_folder, exist_ok=True)

    def process_file(self, input_path: str, output_path: str):
        """
        Processes a single PowerPoint file: converts text to uppercase and saves the result.
        """
        print(f"Processing file: {input_path}")
        editor = PowerPointEditor(input_path)
        editor.convert_text_to_uppercase()
        editor.save(output_path)

    def process_all_files(self):
        """
        Processes all .pptx files in the input folder concurrently and saves the results to the output folder.
        """
        tasks = []
        for filename in filter(lambda f: f.endswith(".pptx"), os.listdir(self.input_folder)):
            input_path = os.path.join(self.input_folder, filename)
            output_filename = os.path.splitext(filename)[0] + "-out.pptx"
            output_path = os.path.join(self.output_folder, output_filename)
            tasks.append((input_path, output_path))

        # Process files concurrently
        with ThreadPoolExecutor() as executor:
            executor.map(lambda task: self.process_file(*task), tasks)


# Usage
# Usage
def main():
    input_folder = "input"  # Folder containing input .pptx files
    output_folder = "output"  # Folder to save the modified .pptx files

    # Create a batch processor and process all files
    batch_processor = PowerPointBatchProcessor(input_folder, output_folder)
    batch_processor.process_all_files()



if __name__ == "__main__":
    main()
